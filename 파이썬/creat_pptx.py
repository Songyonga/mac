from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 새 프레젠테이션 생성
prs = Presentation()

# 슬라이드 1
slide_1 = prs.slides.add_slide(prs.slide_layouts[5])
title_1 = slide_1.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
title_1.text = "취옹정기(醉翁亭) - 구양수(歐陽修)"
title_1.text_frame.paragraphs[0].font.size = Pt(44)
title_1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER  # 가운데 정렬

# 슬라이드 2
slide_2 = prs.slides.add_slide(prs.slide_layouts[5])
content_2 = slide_2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(6))
content_2.text = "구양수(歐陽修 1007-1072)는 자는 영숙(永叔), 호는 취옹(醉翁), 육일거사(六一居士)로, 시호(諡號)는 문충(文忠)으로 송(宋)나라의 정치가(政治家)이며 문인(文人)으로 당송팔대가(唐宋八大家)의 한 사람이었으며 스촨성(四川省) 몐양(綿陽)에서 출생(出生)하여 쓰촨성(四川省)의 지방관(地方官)이었던 부친(父親) 구양관(歐陽觀)이 4세에 타계(他界)하여 어머니를 따라 백부(伯父) 구양엽(歐陽曄)이 추관(推官)으로 있는 수주(隨州, 후베이)에서 기거(寄居)하였다."
content_2.text_frame.paragraphs[0].font.size = Pt(20)

# 슬라이드 3-8 공통 서식과 내용
slides_data = [
    ("〈醉翁亭記〉\n環滁는 皆山也라 其西南諸峰이 林壑尤美하여 望之蔚然而深秀者는 瑯琊也요 山行六七里에 漸聞水聲潺潺而瀉出于兩峰之間者는 釀泉也요 峰回路轉에 有亭翼然하여 臨于泉上者는 醉翁亭也라",
     "*蔚然(위연): 초목이 무성하게 우거져 있는 모양.\n*瑯耶(낭야): 저주에 있는 산 이름.\n*潺潺(잔잔): 물이 졸졸 흐르는 모양\n*釀泉(양천): 샘 이름. 釀은 술을 빚는다는 뜻이다. 이 샘물로 술을 빚으면 술맛이 좋다 하여 양천이라 했다.\n*轉(전): 길이 구불구불하다.\n*翼然(익연): 새가 날개를 활짝 펼친 듯한 모양."),
    
    ("作亭者誰오 山之僧智僊也요 名之者誰오 太守自謂也라 太守與客으로 來飮于此할새 飮少輒醉하고 而年又最高라 故로 自號曰 醉翁也라하니 醉翁之意는 不在酒하고 在乎山水之間也니 山水之樂을 得之心而寓之酒也라",
     "*智仙(지선) : 낭야산의 스님 이름.\n*自謂(자위) : 자기 자신을 뜻한다. 스스로를 말한다.\n*飮少輒醉(음소첨취) : 조금 마셔도 금방 취하다. 輒(첩)은 그러한 때마다\n*醉翁(취옹) : 술에 취한 늙은 이. 구양수의 호다.\n*寓之酒(우지주) : 술을 구실 삼다. 술로 표현하다."),
    
    ("若夫日出而林霏開하고 雲歸而巖穴暝하여 晦明變化者는 山間之朝暮也요 野芳發而幽香하고 嘉木秀而繁陰하며 風霜高潔하고 水落而石出者는 山間之四時也라 朝而往하고 暮而歸에 四時之景不同而樂亦無窮也라",
     "*若夫(약부) : 화제 전환의 連辭. 해석하지 않음이 자연스럽다.\n*日出而林霏開: 해가 뜨면 숲 속의 자욱한 아침 안개 걷히다. 林霏는 숲에 엉킨 안개 같은 작은 물방울.\n*雲歸而巖穴暝: 저녁 구름이 돌아오면 바위 동굴에 어둠이 내리다.\n*晦明(회명): 어두워졌다 밝아졌다 함.\n*野芳(야방): 들에 핀 이름 모를 꽃과 풀.\n*幽香(유향): 그윽한 향기.\n*風霜高潔(풍상고결): 바람이 높이 불고 서리는 맑고 깨끗하다. 높고 푸른 하늘. 시원한 바람.\n*水落而石出: 강물이 줄어드니 강바닥의 돌이 드러나다. 강의 겨울 풍경을 나타낸다.\n*朝而往暮而歸: 아침에 가서 아침 경치를 보고 저녁에 돌아오며 저녁 경치를 구경하다."),
    
    ("至於負者는 歌于塗하고 行者休于樹하며 前者呼하고 後者應하여 傴僂提携하여 往來而不絶者는 滁人遊也요 臨谿而漁하니 谿深而魚肥하고 釀泉爲酒하니 泉冽而酒香이라 山肴野蔌을 雜然而前陳者는 太守宴也니 宴酣之樂은 非絲非竹이라",
     "*負者歌于塗: 짐을 진 사람은 길을 가며 노래 부르다. 塗는 道外 같이 쓰인다.\n*傴僂提携(구루제휴): 허리를 굽혀 손을 잡다. 노인과 어린이로 해석하는 학자도 있다.\n*山肴野蔌(산효야속): 산나물 안주와 푸성귀.\n*宴酣之樂(연감): 잔치가 무르익다.\n*非絲非竹: 현악기도 아니고 관악기도 아니다."),
    
    ("射者中하고 奕者勝하여 觥籌交錯하여 起坐而諠譁者는 衆賓歡也요 蒼顔白髮이 頹乎其間者는 太守醉也라 已而요 夕陽在山하고 人影散亂은 太守歸而賓客從也요 樹林陰翳하여 鳴聲上下는 遊人去而禽鳥樂也라",
     "*觥籌交錯(광주 교착): 벌주 잔과 산가지가 뒤섞여 있다.\n*諠譁(훤화): 왁자지껄하게 떠들썩하다.\n*蒼顔白髮(창안백발): 푸른색을 띤 얼굴과 흰머리.\n*頹乎其間(퇴호기간): 술에 취해 여러 빈객들 사이로 쓰러졌다.\n*已而(이이): 얼마 후.\n*陰翳(음예): 그늘져 어두워지다."),
    
    ("然而禽鳥는 知山林之樂하고 而不知人之樂하며 人은 知從太守遊而樂하고 而不知太守之樂其樂也라 醉能同其樂하고 醒能述以文者는 太守也니 太守는 謂誰오 廬陵歐陽脩也니라",
     "*太守之樂其樂: 태수는 사람들이 즐거워하는 것을 즐거워하다.\n*醉能同其樂: 술에 취해서는 그 즐거움을 함께할 수 있다.\n*醒能述以文者 太守: 깨어나 문장으로 서술하는 자는 태수다.\n*廬陵(여릉): 구양수의 고향으로 강서성 길주.")
]

for content_text, annotation_text in slides_data:
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # 내용 텍스트 상자 추가
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    content_box.text = content_text
    for paragraph in content_box.text_frame.paragraphs:
        paragraph.font.size = Pt(20)

    # 주석 텍스트 상자 추가
    annotation_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
    annotation_box.text = annotation_text
    for paragraph in annotation_box.text_frame.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(0, 0, 255)  # 파란색 주석

# 프레젠테이션 저장
prs.save("/Users/song-yong-a/Desktop/취옹정기_프레젠테이션.pptx")
